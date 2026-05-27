#!/usr/bin/env python3
# /// script
# requires-python = ">=3.11"
# dependencies = ["python-pptx>=1.0.2"]
# ///
"""create_pptx — build a PowerPoint (.pptx) deck from a JSON spec.

Two slide styles, mixable in one deck:

A) STRUCTURED slides — title + bullets/body, laid out by PowerPoint's templates.
B) CUSTOM slides — absolutely-positioned elements (text boxes, rectangles,
   images) with explicit x/y/w/h in INCHES, the same control pptxgenjs gives.

Top-level spec:
{
  "layout":   "WIDE",                  # "WIDE" 13.33x7.5 (16:9, default) or "STANDARD" 10x7.5 (4:3)
  "title":    "Launch Plan",           # optional auto title slide
  "subtitle": "2026 Roadmap",          # optional title-slide subtitle
  "accent":   "1F4E35",                # optional, reserved for future theming
  "slides": [ ...structured or custom slides... ]
}

A) Structured slide:
  { "title": "Goals",
    "bullets": ["Ship v1", {"text": "sub point", "level": 1}],   # strings or {text, level}
    "body":  "Free-form paragraph instead of bullets.",          # optional
    "notes": "Speaker notes." }

B) Custom slide (presence of `elements` switches to absolute positioning):
  { "background": "1F4E35",            # optional slide fill hex
    "notes": "...",
    "elements": [
      { "type": "rect",  "x": 0, "y": 0, "w": 13.33, "h": 1.6, "fill": "1F4E35", "line": "CCCCCC" },
      { "type": "text",  "text": "NICOLE", "x": 0.6, "y": 0.4, "w": 9, "h": 1,
        "size": 30, "bold": true, "color": "FFFFFF", "align": "left", "valign": "middle",
        "font": "Calibri", "fill": "2E7D52" },
      { "type": "bullets", "items": ["Coffee agronomy", "Climate-resilient farming"],
        "x": 0.6, "y": 2.0, "w": 9, "h": 3, "size": 14, "color": "222222", "bullet_char": "•" },
      { "type": "image", "path": "/workspace/work/scratch/photo.png", "x": 10.5, "y": 0.5, "w": 2.2 }
    ]
  }

Element units are INCHES. Colors are 6-digit hex without `#`. `align` ∈
left|center|right; `valign` ∈ top|middle|bottom. For images, give `w` and/or `h`
(omit one to keep aspect ratio).

Examples:
  uv run create_pptx.py --spec /workspace/work/scratch/deck.json --output /workspace/work/exports/deck.pptx
  echo '{"title":"Hi","slides":[{"title":"One","bullets":["a","b"]}]}' | uv run create_pptx.py --stdin -o /workspace/work/exports/d.pptx

Output: single-line JSON -> {"output_path": "...", "format": "pptx", "bytes": N, "slides": 3}
Exit codes: 0 ok | 2 usage/spec error | 3 render error
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path


def _die(msg: str, code: int = 2):
    print(f"create_pptx: {msg}", file=sys.stderr)
    raise SystemExit(code)


def _load_spec(args: argparse.Namespace) -> dict:
    if args.stdin:
        raw = sys.stdin.read()
    elif args.spec_json is not None:
        raw = args.spec_json
    elif args.spec is not None:
        try:
            raw = Path(args.spec).read_text(encoding="utf-8")
        except OSError as e:
            _die(f"cannot read --spec file: {e}")
    else:
        _die("provide one of --spec FILE, --spec-json STR, or --stdin")
    try:
        spec = json.loads(raw)
    except json.JSONDecodeError as e:
        _die(f"spec is not valid JSON: {e}")
    if not isinstance(spec, dict):
        _die("spec must be a JSON object")
    return spec


def _color(value):
    from pptx.dml.color import RGBColor

    v = str(value).strip().lstrip("#")
    if len(v) != 6:
        _die(f"color must be 6-digit hex, got {value!r}")
    try:
        return RGBColor.from_string(v.upper())
    except ValueError:
        _die(f"invalid hex color: {value!r}")


def _set_background(slide, hex_color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = _color(hex_color)


def _render_text_element(slide, el: dict):
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(
        Inches(float(el.get("x", 0))), Inches(float(el.get("y", 0))),
        Inches(float(el.get("w", 4))), Inches(float(el.get("h", 1))),
    )
    if el.get("fill"):
        box.fill.solid()
        box.fill.fore_color.rgb = _color(el["fill"])
    else:
        box.fill.background()
    box.line.fill.background()

    tf = box.text_frame
    tf.word_wrap = True
    valign = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}
    tf.vertical_anchor = valign.get(str(el.get("valign", "top")).lower(), MSO_ANCHOR.TOP)

    p = tf.paragraphs[0]
    align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
    p.alignment = align.get(str(el.get("align", "left")).lower(), PP_ALIGN.LEFT)
    run = p.add_run()
    run.text = str(el.get("text", ""))
    font = run.font
    font.size = Pt(float(el.get("size", 18)))
    font.bold = bool(el.get("bold", False))
    font.italic = bool(el.get("italic", False))
    if el.get("font"):
        font.name = str(el["font"])
    font.color.rgb = _color(el.get("color", "222222"))


def _render_bullets_element(slide, el: dict):
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    box = slide.shapes.add_textbox(
        Inches(float(el.get("x", 0))), Inches(float(el.get("y", 0))),
        Inches(float(el.get("w", 6))), Inches(float(el.get("h", 3))),
    )
    box.fill.background()
    box.line.fill.background()
    tf = box.text_frame
    tf.word_wrap = True
    bullet_char = el.get("bullet_char", "•")
    items = el.get("items", []) or []
    for i, item in enumerate(items):
        text = item if isinstance(item, str) else str(item.get("text", ""))
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet_char} {text}" if bullet_char else text
        run.font.size = Pt(float(el.get("size", 14)))
        run.font.color.rgb = _color(el.get("color", "222222"))
        if el.get("font"):
            run.font.name = str(el["font"])


def _render_rect_element(slide, el: dict):
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(float(el.get("x", 0))), Inches(float(el.get("y", 0))),
        Inches(float(el.get("w", 1))), Inches(float(el.get("h", 1))),
    )
    if el.get("fill"):
        shape.fill.solid()
        shape.fill.fore_color.rgb = _color(el["fill"])
    else:
        shape.fill.background()
    if el.get("line"):
        shape.line.color.rgb = _color(el["line"])
        shape.line.width = Pt(float(el.get("line_width", 1)))
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False


def _render_image_element(slide, el: dict):
    from pptx.util import Inches

    path = el.get("path")
    if not path or not Path(path).exists():
        _die(f"image element path not found: {path!r}", code=3)
    kwargs = {}
    if el.get("w") is not None:
        kwargs["width"] = Inches(float(el["w"]))
    if el.get("h") is not None:
        kwargs["height"] = Inches(float(el["h"]))
    slide.shapes.add_picture(str(path), Inches(float(el.get("x", 0))), Inches(float(el.get("y", 0))), **kwargs)


def _render_custom_slide(prs, sdef: dict):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    if sdef.get("background"):
        _set_background(slide, sdef["background"])
    for el in sdef.get("elements", []) or []:
        if not isinstance(el, dict):
            _die("each element must be an object")
        etype = str(el.get("type", "text")).lower()
        if etype == "text":
            _render_text_element(slide, el)
        elif etype == "bullets":
            _render_bullets_element(slide, el)
        elif etype == "rect":
            _render_rect_element(slide, el)
        elif etype == "image":
            _render_image_element(slide, el)
        else:
            _die(f"unknown element type {etype!r} (use text|bullets|rect|image)")
    if sdef.get("notes"):
        slide.notes_slide.notes_text_frame.text = str(sdef["notes"])
    return slide


def _render_structured_slide(prs, sdef: dict):
    content_layout = prs.slide_layouts[1]
    blank_layout = prs.slide_layouts[6]
    bullets = sdef.get("bullets") or []
    body = sdef.get("body")
    use_content = bool(bullets or body)
    slide = prs.slides.add_slide(content_layout if use_content else blank_layout)

    if slide.shapes.title is not None:
        slide.shapes.title.text = str(sdef.get("title", ""))

    if use_content and len(slide.placeholders) > 1:
        tf = slide.placeholders[1].text_frame
        tf.clear()
        first = True
        if body:
            tf.paragraphs[0].text = str(body)
            first = False
        for b in bullets:
            if isinstance(b, dict):
                text = str(b.get("text", ""))
                level = int(b.get("level", 0) or 0)
            else:
                text = str(b)
                level = 0
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.text = text
            p.level = max(0, min(4, level))

    if sdef.get("notes"):
        slide.notes_slide.notes_text_frame.text = str(sdef["notes"])
    return slide


def _render(spec: dict, output: Path) -> dict:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    layout = str(spec.get("layout", "WIDE")).upper()
    if layout == "WIDE":
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    elif layout == "STANDARD":
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
    else:
        _die(f"layout must be WIDE or STANDARD, got {layout!r}")

    if spec.get("title"):
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = str(spec["title"])
        if spec.get("subtitle") and len(title_slide.placeholders) > 1:
            title_slide.placeholders[1].text = str(spec["subtitle"])

    slides = spec.get("slides", [])
    if not isinstance(slides, list):
        _die("`slides` must be an array")
    for idx, sdef in enumerate(slides):
        if not isinstance(sdef, dict):
            _die(f"slides[{idx}] must be an object")
        if sdef.get("elements"):
            _render_custom_slide(prs, sdef)
        else:
            _render_structured_slide(prs, sdef)

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return {"slides": len(prs.slides._sldIdLst)}


def main(argv: "list[str] | None" = None) -> int:
    ap = argparse.ArgumentParser(
        prog="create_pptx",
        description="Build a .pptx deck from a JSON spec — structured slides and/or absolute-positioned custom slides.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=__doc__,
    )
    src = ap.add_mutually_exclusive_group()
    src.add_argument("--spec", metavar="FILE", help="Path to a JSON spec file.")
    src.add_argument("--spec-json", metavar="STR", help="Inline JSON spec string.")
    src.add_argument("--stdin", action="store_true", help="Read the JSON spec from stdin.")
    ap.add_argument("-o", "--output", required=True, metavar="PATH", help="Output .pptx path (parent dirs auto-created).")
    args = ap.parse_args(argv)

    output = Path(args.output)
    if output.suffix.lower() != ".pptx":
        output = output.with_suffix(".pptx")

    spec = _load_spec(args)
    extra = _render(spec, output)
    print(json.dumps({
        "output_path": str(output.resolve()),
        "format": "pptx",
        "bytes": output.stat().st_size if output.exists() else 0,
        **extra,
    }))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
